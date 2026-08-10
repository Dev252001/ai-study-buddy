"""
PPTX text extractor using python-pptx.
Extracts text from slide shapes, notes, and table cells.
"""
from __future__ import annotations

import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class PPTXProcessor:
    def extract_text(self, file_path: str) -> tuple[str, int]:
        """
        Return *(text, slide_count)*.
        """
        from pptx import Presentation
        from pptx.util import Pt

        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PPTX not found: {file_path}")

        prs = Presentation(file_path)
        slide_texts: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            parts: list[str] = [f"--- Slide {slide_num} ---"]

            for shape in slide.shapes:
                # Text frames (text boxes, titles, content)
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if para_text:
                            parts.append(para_text)

                # Tables
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    parts.append(f"[Notes]: {notes_text}")

            slide_texts.append("\n".join(parts))

        text = "\n\n".join(slide_texts)
        return text, len(prs.slides)
